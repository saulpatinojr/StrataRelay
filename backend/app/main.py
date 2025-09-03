from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import pandas as pd
from google.cloud import storage, firestore, secretmanager
from docx import Document
from pptx import Presentation
import requests
import json
import os

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize clients
storage_client = storage.Client()
db = firestore.Client()
secret_client = secretmanager.SecretManagerServiceClient()

def get_secret(secret_name):
    name = f"projects/{os.getenv('GOOGLE_CLOUD_PROJECT')}/secrets/{secret_name}/versions/latest"
    response = secret_client.access_secret_version(request={"name": name})
    return response.payload.data.decode("UTF-8")

@app.post("/process-file")
async def process_file(job_id: str, file_path: str):
    try:
        # Update job status
        db.collection('jobs').document(job_id).update({'status': 'processing'})
        
        # Download file from storage
        bucket = storage_client.bucket(os.getenv('STORAGE_BUCKET'))
        blob = bucket.blob(file_path)
        file_content = blob.download_as_bytes()
        
        # Process with pandas
        df = pd.read_excel(file_content)
        
        # Generate reports
        doc_url = generate_docx_report(df, job_id)
        ppt_url = generate_pptx_report(df, job_id)
        
        # Push to Power BI
        await push_to_powerbi(df)
        
        # Update job with completion
        db.collection('jobs').document(job_id).update({
            'status': 'completed',
            'reportUrl': doc_url,
            'pptUrl': ppt_url
        })
        
        return {"status": "success", "job_id": job_id}
        
    except Exception as e:
        db.collection('jobs').document(job_id).update({'status': 'error', 'error': str(e)})
        raise HTTPException(status_code=500, detail=str(e))

def generate_docx_report(df, job_id):
    doc = Document()
    doc.add_heading('Analytics Report', 0)
    doc.add_paragraph(f'Data Summary: {len(df)} rows processed')
    
    # Add table
    table = doc.add_table(rows=1, cols=len(df.columns))
    hdr_cells = table.rows[0].cells
    for i, col in enumerate(df.columns):
        hdr_cells[i].text = col
    
    for _, row in df.head(10).iterrows():
        row_cells = table.add_row().cells
        for i, value in enumerate(row):
            row_cells[i].text = str(value)
    
    # Upload to storage
    bucket = storage_client.bucket(os.getenv('STORAGE_BUCKET'))
    blob = bucket.blob(f'reports/{job_id}_report.docx')
    doc.save(f'/tmp/{job_id}_report.docx')
    blob.upload_from_filename(f'/tmp/{job_id}_report.docx')
    
    return blob.public_url

def generate_pptx_report(df, job_id):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Analytics Report"
    subtitle.text = f"Data processed: {len(df)} rows"
    
    # Upload to storage
    bucket = storage_client.bucket(os.getenv('STORAGE_BUCKET'))
    blob = bucket.blob(f'reports/{job_id}_presentation.pptx')
    prs.save(f'/tmp/{job_id}_presentation.pptx')
    blob.upload_from_filename(f'/tmp/{job_id}_presentation.pptx')
    
    return blob.public_url

async def push_to_powerbi(df):
    try:
        powerbi_token = get_secret('powerbi-token')
        dataset_id = get_secret('powerbi-dataset-id')
        
        headers = {
            'Authorization': f'Bearer {powerbi_token}',
            'Content-Type': 'application/json'
        }
        
        data = df.to_dict('records')
        response = requests.post(
            f'https://api.powerbi.com/v1.0/myorg/datasets/{dataset_id}/rows',
            headers=headers,
            json={'rows': data}
        )
        response.raise_for_status()
    except Exception as e:
        print(f"Power BI push failed: {e}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8080)